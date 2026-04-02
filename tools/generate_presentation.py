"""
GraphoLab — Generatore di presentazione PowerPoint
Genera docs/GraphoLab_Presentazione.pptx (~46 slide, italiano)

Uso:
    pip install python-pptx
    python tools/generate_presentation.py
"""

import atexit
import shutil
import tempfile
from pathlib import Path

import matplotlib
matplotlib.use("Agg")  # non-interactive backend
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Palette e costanti
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
GOLD   = RGBColor(0xC8, 0x97, 0x3A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

TITLE_FONT = "Calibri Light"
BODY_FONT  = "Calibri"

OUTPUT  = Path(__file__).parent.parent / "docs" / "GraphoLab_Presentazione.pptx"
SAMPLES = Path(__file__).parent.parent / "data" / "samples"

# Temporary directory for matplotlib PNGs (auto-deleted on exit)
_TMP = tempfile.mkdtemp(prefix="grapholab_pptx_")
atexit.register(shutil.rmtree, _TMP, ignore_errors=True)

# Matplotlib color tuples (0-1 scale)
_NAVY_M  = (0x1B/255, 0x3A/255, 0x6B/255)
_GOLD_M  = (0xC8/255, 0x97/255, 0x3A/255)
_GREEN_M = (0x1A/255, 0x5C/255, 0x38/255)
_RED_M   = (0x8B/255, 0x1A/255, 0x1A/255)
_PURP_M  = (0x5C/255, 0x1A/255, 0x7A/255)
_LBLU_M  = (0x4A/255, 0x7D/255, 0xB8/255)


# ---------------------------------------------------------------------------
# Matplotlib diagram generators
# ---------------------------------------------------------------------------

def _mbox(ax, cx, cy, w, h, text, fc, tc="white", fs=9, bold=True):
    """Draw a rounded box centered at (cx, cy) with text."""
    patch = mpatches.FancyBboxPatch(
        (cx - w / 2, cy - h / 2), w, h,
        boxstyle="round,pad=0.08",
        facecolor=fc, edgecolor="white", linewidth=1.5,
        transform=ax.transData, clip_on=False,
    )
    ax.add_patch(patch)
    ax.text(cx, cy, text, ha="center", va="center",
            fontsize=fs, color=tc, fontweight="bold" if bold else "normal",
            multialignment="center")


def _marrow(ax, x1, y1, x2, y2, color=_GOLD_M, lw=2.0):
    """Draw an arrow from (x1,y1) to (x2,y2)."""
    ax.annotate(
        "", xy=(x2, y2), xytext=(x1, y1),
        arrowprops=dict(arrowstyle="-|>", color=color, lw=lw,
                        mutation_scale=14),
    )


def _gen_training_cycle(path: Path):
    fig, ax = plt.subplots(figsize=(5.5, 4.8), facecolor="white")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")

    # 4 boxes: top, right, bottom, left
    boxes = [
        (5.0, 8.0, "1. Dati\netichettati",   _NAVY_M),
        (8.5, 5.0, "2. Previsione\nmodello",  _LBLU_M),
        (5.0, 2.0, "3. Calcolo\nerrore",      _RED_M),
        (1.5, 5.0, "4. Aggiorna\nparametri",  _PURP_M),
    ]
    for cx, cy, text, fc in boxes:
        _mbox(ax, cx, cy, 2.8, 1.4, text, fc, fs=9)

    # Clockwise arrows (from edge of each box to edge of next)
    pairs = [(0, 1), (1, 2), (2, 3), (3, 0)]
    midpoints = [(b[0], b[1]) for b in boxes]
    for i, j in pairs:
        x1, y1 = midpoints[i]
        x2, y2 = midpoints[j]
        # Move arrow start/end toward the neighbor
        dx, dy = x2 - x1, y2 - y1
        dist = (dx**2 + dy**2) ** 0.5
        off = 1.5
        _marrow(ax, x1 + dx/dist*off, y1 + dy/dist*off,
                    x2 - dx/dist*off, y2 - dy/dist*off, color=_GOLD_M)

    # Center label
    ax.text(5.0, 5.0, "Training\nLoop", ha="center", va="center",
            fontsize=11, color=_NAVY_M, fontweight="bold",
            bbox=dict(boxstyle="circle,pad=0.4", facecolor=(0.95, 0.97, 1.0),
                      edgecolor=_GOLD_M, linewidth=2))

    plt.tight_layout(pad=0.2)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _gen_neural_network(path: Path):
    fig, ax = plt.subplots(figsize=(6.5, 5.0), facecolor="white")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")

    layers = [
        (1.5, [3.0, 5.0, 7.0], _NAVY_M,  "Strato di\ninput"),
        (5.0, [2.0, 4.0, 6.0, 8.0], _GREEN_M, "Strato\nnascosto"),
        (8.5, [3.5, 6.5],    _GOLD_M,  "Strato di\noutput"),
    ]

    # Draw connections first (behind nodes)
    for (x1, ys1, _, _l1), (x2, ys2, _, _l2) in zip(layers, layers[1:]):
        for y1 in ys1:
            for y2 in ys2:
                ax.plot([x1, x2], [y1, y2], color=(0.75, 0.75, 0.85),
                        lw=0.8, zorder=1)

    # Draw nodes
    r = 0.45
    for x, ys, color, label in layers:
        for y in ys:
            circle = plt.Circle((x, y), r, facecolor=color,
                                 edgecolor="white", linewidth=2, zorder=2)
            ax.add_patch(circle)
        # Layer label at bottom
        ax.text(x, 0.7, label, ha="center", va="center",
                fontsize=9, color=_NAVY_M, fontweight="bold",
                multialignment="center")

    # Input feature labels
    feat = ["Tratto", "Forma", "Pressione"]
    for y, name in zip(layers[0][1], feat):
        ax.text(0.1, y, name, ha="left", va="center",
                fontsize=7.5, color=_NAVY_M)

    # Output labels
    out = ["Autentica", "Falsa"]
    for y, name in zip(layers[2][1], out):
        ax.text(9.2, y, name, ha="left", va="center",
                fontsize=8, color=_NAVY_M, fontweight="bold")

    ax.set_title("Rete Neurale — Architettura a Strati", fontsize=10,
                 color=_NAVY_M, fontweight="bold", pad=4)

    plt.tight_layout(pad=0.2)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _gen_computer_vision(path: Path):
    fig, ax = plt.subplots(figsize=(9, 3.2), facecolor="white")
    ax.set_xlim(0, 18)
    ax.set_ylim(0, 4)
    ax.axis("off")

    steps = [
        (1.5, "Immagine\nmanoscritto", _NAVY_M),
        (4.8, "Griglia\npixel",        _LBLU_M),
        (8.1, "Filtri CNN\n(feature)",  _GREEN_M),
        (11.4, "Feature\nmap",          _PURP_M),
        (14.7, "Classificazione\n(output)", _GOLD_M),
    ]

    for x, label, color in steps:
        _mbox(ax, x, 2.0, 2.6, 1.5, label, color, fs=9)

    # Arrows between boxes
    for i in range(len(steps) - 1):
        x1 = steps[i][0] + 1.3
        x2 = steps[i+1][0] - 1.3
        _marrow(ax, x1, 2.0, x2, 2.0, color=_GOLD_M)

    ax.set_title("Pipeline Computer Vision", fontsize=10,
                 color=_NAVY_M, fontweight="bold", pad=4)

    plt.tight_layout(pad=0.2)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _gen_confidence_scores(path: Path):
    fig, axes = plt.subplots(2, 1, figsize=(6, 4.5), facecolor="white")
    fig.subplots_adjust(hspace=0.55)

    examples = [
        ("Esempio 1: Firma Autentica", [91, 9],  [_GREEN_M, _RED_M]),
        ("Esempio 2: Firma Contraffatta", [12, 88], [_GREEN_M, _RED_M]),
    ]
    labels = ["Autentica", "Falsa"]

    for ax, (title, scores, colors) in zip(axes, examples):
        bars = ax.barh(labels, scores, color=colors, edgecolor="white",
                       height=0.55)
        ax.set_xlim(0, 100)
        ax.set_xlabel("Confidenza (%)", fontsize=8, color=_NAVY_M)
        ax.set_title(title, fontsize=9, color=_NAVY_M, fontweight="bold")
        ax.tick_params(labelsize=8)
        ax.spines[["top", "right"]].set_visible(False)
        for bar, score in zip(bars, scores):
            ax.text(min(score + 2, 95), bar.get_y() + bar.get_height()/2,
                    f"{score}%", va="center", fontsize=9, fontweight="bold",
                    color=_NAVY_M)

    plt.tight_layout(pad=0.4)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _gen_siamese_network(path: Path):
    fig, ax = plt.subplots(figsize=(6.5, 5.5), facecolor="white")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")

    # Input boxes
    _mbox(ax, 2.0, 9.0, 2.8, 1.0, "Firma di\nriferimento", _NAVY_M, fs=8)
    _mbox(ax, 8.0, 9.0, 2.8, 1.0, "Firma in\nesame", _NAVY_M, fs=8)

    # CNN boxes (shared weights)
    _mbox(ax, 2.0, 6.8, 2.8, 1.0, "CNN SigNet", _LBLU_M, fs=8)
    _mbox(ax, 8.0, 6.8, 2.8, 1.0, "CNN SigNet\n(stessi pesi)", _LBLU_M, fs=8)

    # Embedding circles
    _mbox(ax, 2.0, 4.6, 2.8, 0.9, "Embedding A", _GREEN_M, fs=8)
    _mbox(ax, 8.0, 4.6, 2.8, 0.9, "Embedding B", _GREEN_M, fs=8)

    # Distance box
    _mbox(ax, 5.0, 2.8, 3.2, 0.9, "Distanza coseno", _GOLD_M, tc=_NAVY_M, fs=9)

    # Verdict box
    _mbox(ax, 5.0, 1.0, 4.0, 1.0,
          "AUTENTICA / FALSA\n(soglia = 0.35)", _RED_M, fs=8)

    # Arrows
    for x in [2.0, 8.0]:
        _marrow(ax, x, 8.5, x, 7.35)   # input → CNN
        _marrow(ax, x, 6.3, x, 5.05)   # CNN → embedding
        _marrow(ax, x, 4.15, 5.0, 3.25)  # embedding → distance

    _marrow(ax, 5.0, 2.35, 5.0, 1.5)   # distance → verdict

    # Shared weights label
    ax.text(5.0, 6.8, "⟵ pesi condivisi ⟶", ha="center", va="center",
            fontsize=7.5, color=_GOLD_M, fontstyle="italic")

    ax.set_title("Architettura Rete Siamese (SigNet)", fontsize=10,
                 color=_NAVY_M, fontweight="bold", pad=4)

    plt.tight_layout(pad=0.2)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _gen_detr_detection(path: Path):
    fig, ax = plt.subplots(figsize=(5.5, 7.0), facecolor="white")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 14)
    ax.axis("off")

    # Document background
    doc = mpatches.FancyBboxPatch((0.5, 0.5), 9.0, 13.0,
                                   boxstyle="round,pad=0.1",
                                   facecolor=(0.97, 0.97, 0.97),
                                   edgecolor=(0.6, 0.6, 0.6), linewidth=1.5)
    ax.add_patch(doc)
    ax.text(5.0, 13.2, "DOCUMENTO SCANSIONATO", ha="center", va="center",
            fontsize=8, color=(0.5, 0.5, 0.5), fontstyle="italic")

    # Horizontal lines (simulating text)
    for y in [11.5, 10.8, 10.1, 9.4, 8.7, 8.0, 7.3, 6.6]:
        w = np.random.uniform(4.0, 7.5)
        ax.plot([1.2, 1.2 + w], [y, y], color=(0.75, 0.75, 0.75), lw=4,
                solid_capstyle="round")

    # Bounding box 1 — main signature
    bb1 = mpatches.FancyBboxPatch((1.5, 1.8), 4.5, 2.8,
                                   boxstyle="square,pad=0.0",
                                   facecolor=(0.0, 0.8, 0.0, 0.08),
                                   edgecolor=_GREEN_M, linewidth=2.5)
    ax.add_patch(bb1)
    ax.text(1.7, 4.75, "Firma  94%", fontsize=8.5, color=_GREEN_M,
            fontweight="bold",
            bbox=dict(facecolor=_GREEN_M, edgecolor="none", pad=2))
    ax.text(1.6, 3.2, "~~~~~~~~~~~~~~~~", fontsize=14,
            color=(0.3, 0.3, 0.3), fontstyle="italic")

    # Bounding box 2 — second smaller signature
    bb2 = mpatches.FancyBboxPatch((5.8, 3.0), 3.0, 1.6,
                                   boxstyle="square,pad=0.0",
                                   facecolor=(0.0, 0.8, 0.0, 0.08),
                                   edgecolor=_GREEN_M, linewidth=2.5)
    ax.add_patch(bb2)
    ax.text(5.9, 4.72, "Firma  87%", fontsize=8.5, color=_GREEN_M,
            fontweight="bold",
            bbox=dict(facecolor=_GREEN_M, edgecolor="none", pad=2))
    ax.text(5.9, 3.7, "~~~~~~~~~", fontsize=14,
            color=(0.3, 0.3, 0.3), fontstyle="italic")

    ax.set_title("Conditional DETR — Rilevamento firme nel documento", fontsize=9,
                 color=_NAVY_M, fontweight="bold", pad=4)

    plt.tight_layout(pad=0.2)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _gen_pipeline_e2e(path: Path):
    fig, ax = plt.subplots(figsize=(10.5, 4.5), facecolor="white")
    ax.set_xlim(0, 21)
    ax.set_ylim(0, 8)
    ax.axis("off")

    # Row 1: document → DETR → SigNet → Autenticità
    r1_y = 6.0
    r1 = [
        (1.5,  r1_y, "Documento\nscansionato", _NAVY_M),
        (5.0,  r1_y, "DETR\nRilevamento",    _LBLU_M),
        (8.5,  r1_y, "Firma\nestratta",         _PURP_M),
        (12.0, r1_y, "SigNet\nVerifica",        _GREEN_M),
        (16.0, r1_y, "Verdetto\nautenticità",   _GOLD_M),
    ]
    for x, y, label, fc in r1:
        _mbox(ax, x, y, 2.6, 1.2, label, fc, fs=8)
    for i in range(len(r1)-1):
        _marrow(ax, r1[i][0]+1.3, r1[i][1], r1[i+1][0]-1.3, r1[i+1][1])

    # Row 2: TrOCR → NER → Entità
    r2_y = 2.8
    r2 = [
        (5.0,  r2_y, "TrOCR\nTrascrizione", _LBLU_M),
        (8.5,  r2_y, "Testo\ntrascritto",    _NAVY_M),
        (12.0, r2_y, "WikiNEural\nNER",      _PURP_M),
        (16.0, r2_y, "Entità\n(persone, luoghi)", _GREEN_M),
    ]
    for x, y, label, fc in r2:
        _mbox(ax, x, y, 2.6, 1.2, label, fc, fs=8)
    for i in range(len(r2)-1):
        _marrow(ax, r2[i][0]+1.3, r2[i][1], r2[i+1][0]-1.3, r2[i+1][1])

    # Vertical arrow from documento (row1) down to TrOCR row
    _marrow(ax, r1[0][0], r1_y - 0.6, r2[0][0], r2_y + 0.6, color=_GOLD_M)

    # Converging arrow to final report
    _mbox(ax, 19.5, 4.4, 2.6, 1.2, "Referto\nforense", _RED_M, fs=8)
    _marrow(ax, r1[-1][0]+1.3, r1[-1][1], 19.5-1.3, 4.9, color=_GOLD_M)
    _marrow(ax, r2[-1][0]+1.3, r2[-1][1], 19.5-1.3, 3.9, color=_GOLD_M)

    ax.set_title("Pipeline End-to-End: dal documento al referto forense",
                 fontsize=10, color=_NAVY_M, fontweight="bold", pad=6)

    plt.tight_layout(pad=0.3)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _gen_signatures_comparison(path: Path):
    """Side-by-side genuine vs forged signature images."""
    genuine_path = SAMPLES / "genuine_1_1.png"
    forged_path  = SAMPLES / "forged_1_1.png"

    if not genuine_path.exists() or not forged_path.exists():
        # Fallback: draw placeholder
        fig, axes = plt.subplots(1, 2, figsize=(7, 3), facecolor="white")
        for ax, label in zip(axes, ["Firma autentica", "Firma contraffatta"]):
            ax.text(0.5, 0.5, label, ha="center", va="center",
                    transform=ax.transAxes, fontsize=12)
            ax.axis("off")
        fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return

    img_g = np.array(Image.open(genuine_path).convert("RGB"))
    img_f = np.array(Image.open(forged_path).convert("RGB"))

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(7.5, 3.5),
                                    facecolor="white")

    for ax, img, label, frame_color, verdict in [
        (ax1, img_g, "Firma di riferimento", _GREEN_M, "AUTENTICA"),
        (ax2, img_f, "Firma in esame",        _RED_M,   "FALSA"),
    ]:
        ax.imshow(img, cmap="gray" if img.ndim == 2 else None)
        ax.axis("off")
        ax.set_title(label, fontsize=9, color=_NAVY_M, fontweight="bold",
                     pad=4)
        for spine in ax.spines.values():
            spine.set_visible(True)
            spine.set_edgecolor(frame_color)
            spine.set_linewidth(4)
        ax.text(0.5, -0.08, verdict, transform=ax.transAxes,
                ha="center", va="top", fontsize=10, fontweight="bold",
                color=frame_color)

    plt.suptitle("Esempio: firma autentica vs contraffatta (dataset CEDAR)",
                 fontsize=9, color=_NAVY_M, y=1.02)
    plt.tight_layout(pad=0.5)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def generate_diagrams() -> dict:
    """Generate all matplotlib PNGs and return {key: Path}."""
    tmp = Path(_TMP)
    imgs = {}

    tasks = [
        ("training_cycle",       _gen_training_cycle),
        ("neural_network",       _gen_neural_network),
        ("computer_vision",      _gen_computer_vision),
        ("confidence_scores",    _gen_confidence_scores),
        ("siamese_network",      _gen_siamese_network),
        ("detr_detection",       _gen_detr_detection),
        ("pipeline_e2e",         _gen_pipeline_e2e),
        ("signatures_comparison", _gen_signatures_comparison),
    ]

    for key, fn in tasks:
        p = tmp / f"{key}.png"
        fn(p)
        imgs[key] = p
        print(f"  [diagram] {key}.png")

    # Real sample images
    for key, filename in [
        ("handwritten", "handwritten_text_01.png"),
        ("multiline",   "handwritten_multiline_01.png"),
    ]:
        p = SAMPLES / filename
        if p.exists():
            imgs[key] = p

    return imgs


# ---------------------------------------------------------------------------
# Slide helper functions
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 font_name=BODY_FONT, font_size=20,
                 color=DARK, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txb


def _add_bullet_textbox(slide, bullets, left, top, width, height,
                        font_size=18, color=DARK, indent_level=0):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        text, level = (item if isinstance(item, tuple) else (item, indent_level))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("• " if level == 0 else "  – ") + text
        run.font.name = BODY_FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txb


def _add_picture_fill(slide, img_path, left, top, max_w, max_h):
    """Add picture preserving aspect ratio within a max_w × max_h box."""
    if not img_path or not Path(img_path).exists():
        return
    img = Image.open(img_path)
    iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    actual_left = left + (max_w - w) // 2
    actual_top  = top  + (max_h - h) // 2
    slide.shapes.add_picture(str(img_path), actual_left, actual_top, w, h)


# ---------------------------------------------------------------------------
# Slide layout builders
# ---------------------------------------------------------------------------

def _set_speaker_note(slide, text: str):
    """Write text to the PowerPoint notes pane (visible in presenter view)."""
    if not text:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = text.strip()


def add_cover(prs: Presentation, speaker_note: str = ""):
    slide = _blank_slide(prs)
    _fill_bg(slide, NAVY)
    _add_rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55), GOLD)
    _add_textbox(
        slide, "Intelligenza Artificiale\ne Grafologia Forense",
        Inches(1), Inches(1.2), Inches(11.3), Inches(2.8),
        font_name=TITLE_FONT, font_size=44, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, "Nuovi strumenti a supporto dell'esaminatore forense",
        Inches(1), Inches(4.1), Inches(11.3), Inches(0.8),
        font_name=BODY_FONT, font_size=24, color=GOLD, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, "GraphoLab",
        Inches(0.4), SLIDE_H - Inches(0.5), Inches(3), Inches(0.45),
        font_name=TITLE_FONT, font_size=16, color=DARK, bold=True,
    )
    _add_textbox(
        slide, "Presentazione riservata — uso interno",
        Inches(0.4), Inches(6.4), Inches(11), Inches(0.4),
        font_size=13, color=RGBColor(0xAA, 0xBB, 0xCC),
        align=PP_ALIGN.CENTER,
    )
    _set_speaker_note(slide, speaker_note)


def add_agenda(prs: Presentation, speaker_note: str = ""):
    slide = _blank_slide(prs)
    _fill_bg(slide, LIGHT)
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, GOLD)
    _add_textbox(slide, "Agenda", Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
                 font_name=TITLE_FONT, font_size=32, color=NAVY, bold=True)
    items = [
        "1.  Cos'è l'Intelligenza Artificiale",
        "2.  Limiti dell'analisi manuale e ruolo dell'AI",
        "3.  Sei aree di applicazione alla grafologia forense",
        "4.  GraphoLab — i laboratori dimostrativi",
        "5.  Considerazioni etiche e conclusioni",
    ]
    _add_bullet_textbox(slide, items, Inches(0.9), Inches(1.4), Inches(11.5),
                        Inches(5), font_size=22, color=DARK)
    _set_speaker_note(slide, speaker_note)


def add_section(prs: Presentation, number: str, title: str, subtitle: str = "",
                speaker_note: str = ""):
    slide = _blank_slide(prs)
    _fill_bg(slide, NAVY)
    _add_textbox(slide, number, Inches(0.6), Inches(0.5), Inches(2), Inches(1),
                 font_name=TITLE_FONT, font_size=48, color=GOLD, bold=True)
    _add_textbox(slide, title, Inches(0.6), Inches(1.6), Inches(11.5), Inches(2.5),
                 font_name=TITLE_FONT, font_size=40, color=WHITE, bold=True)
    if subtitle:
        _add_textbox(slide, subtitle, Inches(0.6), Inches(4.2), Inches(11.5),
                     Inches(1), font_size=20, color=GOLD, italic=True)
    _add_rect(slide, Inches(0.6), Inches(4.0), Inches(11), Inches(0.06), GOLD)
    _set_speaker_note(slide, speaker_note)


def add_bullet(prs: Presentation, title: str, bullets, note: str = "",
               font_size: int = 19, speaker_note: str = ""):
    slide = _blank_slide(prs)
    _fill_bg(slide, WHITE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    _add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_name=TITLE_FONT, font_size=28, color=WHITE, bold=True)
    _add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), GOLD)
    _add_bullet_textbox(slide, bullets, Inches(0.5), Inches(1.35), Inches(12.3),
                        Inches(5.5), font_size=font_size, color=DARK)
    if note:
        _add_rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55),
                  RGBColor(0xF0, 0xF2, 0xF7))
        _add_textbox(slide, note, Inches(0.4), SLIDE_H - Inches(0.48),
                     Inches(12.5), Inches(0.42), font_size=12, color=GREY, italic=True)
    _set_speaker_note(slide, speaker_note)


def add_bullet_with_image(prs: Presentation, title: str, bullets, img_path,
                          font_size: int = 17, note: str = "", speaker_note: str = ""):
    """Bullet slide with image panel on the right (58% text / 40% image)."""
    slide = _blank_slide(prs)
    _fill_bg(slide, WHITE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    _add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_name=TITLE_FONT, font_size=28, color=WHITE, bold=True)
    _add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), GOLD)

    bullet_w = Inches(7.5)
    _add_bullet_textbox(slide, bullets, Inches(0.5), Inches(1.35), bullet_w,
                        Inches(5.5), font_size=font_size, color=DARK)

    # Image panel on the right
    img_left = Inches(8.15)
    img_top  = Inches(1.35)
    img_w    = Inches(4.9)
    img_h    = Inches(5.5)
    _add_rect(slide, img_left - Inches(0.08), img_top,
              img_w + Inches(0.16), img_h, LIGHT)
    _add_picture_fill(slide, img_path, img_left, img_top + Inches(0.2),
                      img_w, img_h - Inches(0.4))

    if note:
        _add_rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55),
                  RGBColor(0xF0, 0xF2, 0xF7))
        _add_textbox(slide, note, Inches(0.4), SLIDE_H - Inches(0.48),
                     Inches(12.5), Inches(0.42), font_size=12, color=GREY, italic=True)
    _set_speaker_note(slide, speaker_note)


def add_two_col(prs: Presentation, title: str,
                left_title: str, left_bullets,
                right_title: str, right_bullets,
                left_color=NAVY, right_color=RGBColor(0x1A, 0x5C, 0x38),
                left_image=None, speaker_note: str = ""):
    slide = _blank_slide(prs)
    _fill_bg(slide, WHITE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    _add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_name=TITLE_FONT, font_size=28, color=WHITE, bold=True)
    _add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), GOLD)

    col_w = Inches(6.0)
    gap   = Inches(0.35)
    top   = Inches(1.35)
    h     = Inches(5.7)

    # Left column
    _add_rect(slide, Inches(0.25), top, col_w, Inches(0.45), left_color)
    _add_textbox(slide, left_title,
                 Inches(0.35), top + Pt(4), col_w - Inches(0.1), Inches(0.4),
                 font_size=16, color=WHITE, bold=True)

    if left_image and Path(left_image).exists():
        # Reduce bullet area and add image below
        bullet_h = Inches(2.8)
        _add_bullet_textbox(slide, left_bullets,
                            Inches(0.35), top + Inches(0.5), col_w - Inches(0.1),
                            bullet_h, font_size=16, color=DARK)
        img_top_l = top + Inches(0.5) + bullet_h + Inches(0.1)
        img_h_l   = h - Inches(0.5) - bullet_h - Inches(0.2)
        _add_rect(slide, Inches(0.25), img_top_l, col_w, img_h_l, LIGHT)
        _add_picture_fill(slide, left_image,
                          Inches(0.35), img_top_l + Inches(0.1),
                          col_w - Inches(0.2), img_h_l - Inches(0.2))
    else:
        _add_bullet_textbox(slide, left_bullets,
                            Inches(0.35), top + Inches(0.5), col_w - Inches(0.1),
                            h - Inches(0.5), font_size=17, color=DARK)

    # Right column
    rx = Inches(0.25) + col_w + gap
    _add_rect(slide, rx, top, col_w, Inches(0.45), right_color)
    _add_textbox(slide, right_title,
                 rx + Inches(0.1), top + Pt(4), col_w - Inches(0.1), Inches(0.4),
                 font_size=16, color=WHITE, bold=True)
    _add_bullet_textbox(slide, right_bullets,
                        rx + Inches(0.1), top + Inches(0.5), col_w - Inches(0.1),
                        h - Inches(0.5), font_size=17, color=DARK)
    _set_speaker_note(slide, speaker_note)


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list,
                    speaker_note: str = ""):
    slide = _blank_slide(prs)
    _fill_bg(slide, WHITE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    _add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_name=TITLE_FONT, font_size=28, color=WHITE, bold=True)
    _add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), GOLD)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_left   = Inches(0.4)
    tbl_top    = Inches(1.35)
    tbl_width  = Inches(12.5)
    tbl_height = Inches(0.48) * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top,
                                   tbl_width, tbl_height).table
    col_w = int(tbl_width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w

    def _set_cell(cell, text, bg, fg, bold=False, size=15):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = fg
        run.font.bold = bold

    for j, h in enumerate(headers):
        _set_cell(table.cell(0, j), h, NAVY, WHITE, bold=True, size=15)
    for i, row in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            _set_cell(table.cell(i + 1, j), val, bg, DARK, size=14)
    _set_speaker_note(slide, speaker_note)


# ---------------------------------------------------------------------------
# Presentation content
# ---------------------------------------------------------------------------

def build(prs: Presentation, imgs: dict):

    # === COVER + AGENDA ===
    add_cover(prs, speaker_note="""
Buongiorno a tutti e benvenuti. Oggi vedremo come l'intelligenza artificiale può diventare uno strumento potente nelle mani del perito grafologo e dell'investigatore forense. Non si tratta di sostituire la competenza umana, ma di amplificarla con strumenti che lavorano su larga scala e producono misurazioni oggettive. Iniziamo dall'inizio: cos'è davvero l'intelligenza artificiale?
""")
    add_agenda(prs, speaker_note="""
Ecco la struttura della presentazione. Partiamo dalle basi — cos'è l'AI, come funziona — senza dare nulla per scontato. Poi entriamo nel vivo: i limiti dell'analisi manuale, le sei aree di applicazione, i laboratori pratici. Chiudiamo con le considerazioni etiche e la demo interattiva che potete esplorare direttamente.
""")

    # =========================================================
    # SEZIONE 1 — Cos'è l'Intelligenza Artificiale
    # =========================================================
    add_section(prs, "01", "Cos'è l'Intelligenza Artificiale",
                "Una introduzione accessibile per professionisti non tecnici",
                speaker_note="""
Partiamo dalle fondamenta. Questo capitolo è dedicato a chi non ha background tecnico: spiegheremo l'AI con analogie concrete, senza formule matematiche. L'obiettivo è costruire un vocabolario comune che ci servirà per il resto della presentazione.
""")

    add_bullet(prs, "Cos'è l'Intelligenza Artificiale", [
        "L'intelligenza artificiale (AI) è la capacità di un sistema informatico di svolgere compiti che normalmente richiedono intelligenza umana",
        "Non si tratta di un singolo algoritmo, ma di una famiglia di tecniche: machine learning, reti neurali, computer vision, elaborazione del linguaggio naturale",
        "L'AI non «pensa» come un essere umano — riconosce pattern statistici in grandi quantità di dati",
        "Il risultato è un sistema capace di generalizzare: applicare ciò che ha imparato a casi nuovi, mai visti prima",
    ], note="AI = riconoscimento di pattern su larga scala, non intelligenza nel senso umano del termine",
    speaker_note="""
La prima cosa da chiarire: l'AI non è magia, non è fantascienza e non 'pensa' nel senso umano. È un sistema statistico che ha visto milioni di esempi e ha imparato a riconoscere pattern. Pensate a come un bambino impara a riconoscere un gatto: vede mille gatti diversi finché il suo cervello costruisce un'idea generale di 'gatto'. L'AI fa esattamente questo, ma con dati numerici e su scala enormemente maggiore.
""")

    add_bullet(prs, "AI vs. Regole Manuali: Programmare vs. Imparare", [
        "Approccio tradizionale: un programmatore scrive regole esplicite  →  il computer le esegue",
        "  Esempio: «Se il pixel è scuro E la forma è tonda, allora è la lettera O»",
        "Approccio AI (machine learning): il sistema riceve migliaia di esempi etichettati e deduce autonomamente le regole",
        "  Esempio: il modello vede 100.000 lettere «O» scritte a mano e impara da solo come riconoscerle",
        "Vantaggio: i sistemi AI gestiscono variabilità e complessità che sarebbe impossibile codificare manualmente",
    ], speaker_note="""
L'approccio tradizionale è: un programmatore scrive tutte le regole a mano. Ma la scrittura a mano è infinitamente variabile: impossibile scrivere una regola per ogni stile di ogni persona. Il machine learning capovolge l'approccio: invece di scrivere le regole, diamo al sistema migliaia di esempi etichettati e lui le deduce da solo. Questo è il salto concettuale fondamentale dell'AI moderna.
""")

    add_bullet_with_image(prs, "Come Impara una Macchina: il Training",
        [
            "Training: il modello analizza esempi e aggiusta i propri parametri interni",
            "Fase 1 — Input: dati etichettati (es. firme con etichetta «autentica» o «falsa»)",
            "Fase 2 — Previsione: il modello produce un output con score di confidenza",
            "Fase 3 — Errore: si misura quanto la previsione si discosta dall'etichetta reale",
            "Fase 4 — Aggiornamento: i parametri vengono modificati per ridurre l'errore",
            "Il ciclo si ripete su milioni di esempi fino a convergenza",
        ],
        imgs.get("training_cycle"),
        note="Il training avviene una sola volta. L'uso del modello addestrato (inferenza) è istantaneo.",
        speaker_note="""
Il diagramma mostra il ciclo di training. È come un bambino che impara con le flash card: vede l'esempio, fa un tentativo, riceve una correzione, aggiusta. Questo ciclo si ripete milioni di volte. La cosa importante da capire: il training avviene una volta sola — può richiedere ore o giorni su hardware specializzato. Dopo, usare il modello addestrato è istantaneo, come caricare un esperto virtuale che ha già studiato tutto.
""")

    add_bullet(prs, "I Dati di Addestramento: il Carburante dell'AI", [
        "Senza dati non c'è apprendimento: la qualità del modello dipende direttamente dalla qualità e quantità dei dati di training",
        "Dataset tipici per grafologia forense: migliaia di firme etichettate, campioni di scrittura da più scrittori, documenti scansionati",
        "Bias nei dati: se il dataset è sbilanciato, il modello potrebbe funzionare peggio su gruppi sotto-rappresentati",
        "I modelli usati in GraphoLab sono addestrati su dataset pubblici di riferimento: GPDS (firme), IAM (scrittura), CEDAR (firme)",
    ], speaker_note="""
Si dice 'garbage in, garbage out': se il dataset è sbilanciato o di bassa qualità, il modello sarà sbilanciato. Per la grafologia forense questo è critico: un modello addestrato solo su firme di un certo periodo o cultura potrebbe funzionare meno bene su altri. I modelli di GraphoLab usano dataset pubblici di riferimento — GPDS, IAM, CEDAR — che sono gli standard del settore, con migliaia di campioni verificati.
""")

    add_bullet_with_image(prs, "Reti Neurali: Ispirazione Biologica",
        [
            "Le reti neurali artificiali sono ispirate al funzionamento del cervello: unità di calcolo (neuroni) collegate in strati",
            "Ogni connessione ha un peso numerico aggiustato durante il training",
            "Strato di input  →  strati nascosti (elaborazione)  →  strato di output",
            "Le reti imparano rappresentazioni gerarchiche: dai tratti alle lettere, alle parole",
            "Esempio: una rete per le firme impara prima a riconoscere bordi e curve, poi pattern stilistici",
        ],
        imgs.get("neural_network"),
        speaker_note="""
Le reti neurali sono ispirate al cervello, ma non antropomorfizzatele troppo. Sono sistemi matematici con milioni di parametri numerici aggiustati durante il training. Il concetto chiave è la gerarchia: i primi strati riconoscono elementi semplici come bordi, i successivi riconoscono forme più complesse, i finali riconoscono strutture ad alto livello. Per le firme: prima bordi e curve, poi loop e tratti caratteristici, poi lo stile complessivo dello scrittore.
""")

    add_bullet_with_image(prs, "Computer Vision: Vedere con le Macchine",
        [
            "La computer vision insegna alle macchine a interpretare immagini",
            "Un'immagine è una griglia di pixel — valori numerici di intensità luminosa",
            "Le reti convoluzionali (CNN) applicano filtri che rilevano bordi, texture e forme a scale diverse",
            "In grafologia forense, la computer vision è lo strumento per:",
            "  Leggere testo manoscritto (HTR)",
            "  Rilevare e verificare firme",
            "  Misurare caratteristiche grafologiche",
        ],
        imgs.get("computer_vision"),
        speaker_note="""
Per noi un'immagine è qualcosa di visivo e significativo. Per la macchina è una griglia di numeri: valori da 0 a 255 per ogni pixel. I filtri convoluzionali scorrono su questa griglia cercando pattern — bordi, angoli, texture. Questi pattern vengono combinati gerarchicamente fino a produrre una classificazione. In grafologia, questo significa che la macchina può misurare l'angolo di ogni tratto con precisione che nessun occhio umano potrebbe raggiungere in modo consistente.
""")

    add_bullet_with_image(prs, "Output Probabilistici: Score, Non Verdetti",
        [
            "I sistemi AI non producono risposte binarie assolute — producono distribuzioni di probabilità",
            "Esempio firma: «Autentica con confidenza 91%» oppure «Falsa con confidenza 85%»",
            "Una soglia decisionale converte la probabilità in una classificazione",
            "La soglia deve essere calibrata sul tipo di rischio accettabile nel contesto forense",
            "Regola fondamentale: lo score AI è un dato quantitativo — l'interpretazione spetta all'esaminatore",
        ],
        imgs.get("confidence_scores"),
        note="Un sistema AI non 'sa' se una firma è falsa. Sa che la distanza coseno è 0.62, che supera la soglia 0.35.",
        speaker_note="""
Questo è un punto cruciale da comunicare bene, soprattutto in sede giudiziaria. L'AI non dice 'questa firma è falsa'. Dice 'la distanza coseno tra questa firma e il riferimento è 0.62, che supera la soglia calibrata di 0.35'. Tutto il resto — interpretare questo numero nel contesto del caso specifico — spetta all'esaminatore. La soglia può essere regolata a seconda del livello di rischio accettabile: una banca e un tribunale penale potrebbero calibrarla diversamente.
""")

    add_bullet(prs, "AI nella Vita Quotidiana: Esempi Familiari", [
        "Traduzione automatica (Google Translate, DeepL) — comprensione del linguaggio naturale",
        "Filtro antispam delle email — classificazione testi",
        "Riconoscimento vocale (Siri, Alexa) — conversione audio → testo",
        "Face ID sullo smartphone — riconoscimento facciale",
        "Diagnosi medica da immagini radiologiche — computer vision applicata alla medicina",
        "In tutti questi casi: l'AI non «capisce», ma ha imparato a riconoscere pattern statistici con alta affidabilità",
    ], speaker_note="""
Per dare concretezza, ricordiamo che tutti usiamo l'AI ogni giorno senza saperlo: il filtro spam della posta, la traduzione automatica, il riconoscimento facciale del telefono. Questi sono tutti esempi dello stesso principio: riconoscimento di pattern su grandi quantità di dati. La differenza rispetto all'uso forense è che in forensics dobbiamo documentare ogni passo e spiegarlo in modo comprensibile a giudici e giurie.
""")

    add_bullet(prs, "Cosa l'AI Non È", [
        "Non è infallibile: i modelli commettono errori, specialmente su casi anomali o dati degradati",
        "Non è trasparente per default: i modelli deep learning sono «scatole nere»",
        "Non sostituisce l'esperienza: non ha accesso al contesto del caso, alla storia del documento",
        "Non produce certezze legali: uno score al 95% non equivale a prova oltre ogni ragionevole dubbio",
        "Non è neutrale per default: riflette i bias presenti nei dati di addestramento",
        "L'AI è uno strumento potente — come un microscopio: trasforma il lavoro dell'esperto, non lo elimina",
    ], note="Il valore dell'AI in ambito forense dipende dal modo in cui viene integrata nel processo peritale.",
    speaker_note="""
Questo è forse il punto più importante. L'AI non è infallibile — commette errori, specialmente su casi ai margini della distribuzione su cui è stata addestrata. Non è trasparente: i modelli deep learning sono 'scatole nere' e spiegare le loro decisioni richiede strumenti appositi. Non ha accesso al contesto del caso. Non produce certezze legali: un 95% di confidenza non è 'prova oltre ogni ragionevole dubbio'. È uno strumento — potente come un microscopio, ma il microscopio non emette diagnosi. Questo porta naturalmente alla domanda: dove l'analisi manuale mostra i suoi limiti?
""")

    # =========================================================
    # SEZIONE 2 — Limiti del metodo manuale e ruolo dell'AI
    # =========================================================
    add_section(prs, "02", "Limiti dell'Analisi Manuale\ne Ruolo dell'AI",
                "Dove l'AI aggiunge valore reale al lavoro dell'esaminatore forense",
                speaker_note="""
Passiamo alla domanda pratica: perché un perito esperto dovrebbe voler usare l'AI? Non per moda o innovazione — ma perché esistono limiti strutturali dell'analisi puramente manuale che l'AI può superare. Vediamoli uno per uno.
""")

    add_bullet(prs, "Limiti dell'Analisi Manuale: Tempo e Scalabilità", [
        "Il confronto manuale di firme o campioni di scrittura è un processo lento e meticoloso",
        "Un esaminatore esperto può analizzare in dettaglio pochi campioni al giorno",
        "Nelle indagini complesse — contratti multi-pagina, archivi finanziari — il volume supera la capacità manuale",
        "Esempio pratico: screening di 10.000 pagine per individuare firme anomale  →  settimane di lavoro manuale vs. poche ore con AI",
        "L'AI non si stanca, non si distrae e applica gli stessi criteri a ogni documento",
    ], speaker_note="""
Immaginate di dover esaminare 50.000 pagine di documentazione bancaria per trovare firme anomale. Un perito esperto, lavorando con rigore, può analizzare qualche decina di firme al giorno. L'AI può pre-filtrare questo archivio in poche ore, segnalando i casi sospetti per la revisione umana. Il perito si concentra sull'1-2% critico, non sull'altro 98%. Questo non riduce il ruolo del perito — lo rende molto più efficace.
""")

    add_bullet(prs, "Limiti dell'Analisi Manuale: Soggettività e Riproducibilità", [
        "L'analisi visiva è intrinsecamente soggettiva: due periti esperti possono giungere a conclusioni diverse sullo stesso campione borderline",
        "Non si tratta di incompetenza — è una caratteristica strutturale di qualsiasi valutazione basata sul giudizio visivo",
        "I tribunali richiedono metodologie verificabili e riproducibili da terze parti",
        "L'AI produce valori numerici precisi (es. inclinazione media 11.3°, spaziatura 4.2 mm) verificabili indipendentemente",
        "Questo non elimina il giudizio del perito, ma lo ancora a misurazioni oggettive",
    ], speaker_note="""
Non è una critica ai periti — è una caratteristica strutturale del giudizio visivo umano. Due radiologi esperti non vedono esattamente la stessa cosa guardando la stessa radiografia. Due sommelier esperti non danno lo stesso voto allo stesso vino. La differenza in forensics è che i tribunali chiedono metodologie verificabili e riproducibili da terze parti. Quando l'AI dice 'inclinazione media 11.3 gradi', questa misurazione è verificabile indipendentemente da chiunque con lo stesso strumento.
""")

    add_bullet(prs, "Limiti dell'Analisi Manuale: Scala degli Archivi Digitali", [
        "La digitalizzazione ha trasformato le controversie legali: i casi coinvolgono sempre più grandi archivi di documenti scansionati",
        "Esaminare manualmente 50.000 pagine di documentazione bancaria non è fattibile nei tempi processuali",
        "L'AI può pre-filtrare: individuare automaticamente le pagine con firme, segnalare anomalie, produrre priorità",
        "L'esaminatore si concentra poi sui casi segnalati come sospetti — moltiplicando l'efficacia del suo lavoro",
    ], speaker_note="""
I processi moderni — frodi finanziarie, successioni contestate, falsi documentali — coinvolgono sempre più archivi digitali enormi. Non è realistico aspettarsi che l'analisi manuale tenga il passo con il volume di documenti nelle controversie commerciali e finanziarie contemporanee. L'AI cambia il modello operativo: invece di 'esamina tutto manualmente', diventa 'AI pre-filtra, esperto verifica i casi segnalati'. Il risultato è un uso molto più efficiente del tempo del perito.
""")

    add_bullet(prs, "Cosa l'AI Aggiunge: Misurazioni Oggettive e Audit Trail", [
        "Misurazioni quantitative: l'AI produce numeri, non impressioni — angoli, distanze, score di similarità",
        "Riproducibilità: chiunque esegua lo stesso modello sullo stesso file ottiene lo stesso risultato",
        "Audit trail: ogni passo dell'analisi AI può essere registrato e documentato",
        "Il referto peritale supportato da misurazioni AI è più resistente alla contestazione tecnica in sede processuale",
        "Standard internazionali (OSAC, SWGDOC) stanno integrando le linee guida per gli strumenti computazionali",
    ], speaker_note="""
Tre parole chiave: numeri, riproducibilità, documentazione. L'AI non dice 'questa firma sembra diversa'. Dice 'distanza coseno = 0.62'. Chiunque riesegua l'analisi con lo stesso modello e lo stesso file ottiene lo stesso numero — zero margine di interpretazione soggettiva nella misurazione. Ogni passo è loggato con timestamp, versione del modello e parametri. Un referto peritale supportato da questi dati è molto più resistente alla contestazione tecnica in aula.
""")

    add_bullet(prs, "Modello di Collaborazione Uomo-AI", [
        "L'AI non è un sostituto dell'esaminatore forense — è un amplificatore delle sue capacità",
        "Cosa fa l'AI:",
        "  Elabora grandi volumi di documenti in tempi brevi",
        "  Produce misurazioni oggettive e riproducibili",
        "  Segnala anomalie e casi sospetti per la revisione umana",
        "Cosa fa l'esaminatore:",
        "  Interpreta i risultati nel contesto del caso",
        "  Applica conoscenza peritale specializzata",
        "  Produce un'opinione professionale con responsabilità legale",
    ], note="Il modello ideale: AI per la quantificazione, esperto per l'interpretazione.",
    speaker_note="""
Il messaggio centrale è questo: l'AI fa il lavoro quantitativo su larga scala — misurazioni, screening, prioritizzazione. L'esperto porta ciò che la macchina non può: contesto del caso, conoscenza specializzata delle dinamiche della scrittura, giudizio professionale, responsabilità legale. Non è competizione, è specializzazione. Passiamo ora a vedere le sei aree concrete dove questo modello si applica.
""")

    # =========================================================
    # SEZIONE 3 — Sei Aree di Applicazione
    # =========================================================
    add_section(prs, "03", "Sei Aree di Applicazione\nalla Grafologia Forense",
                "Da HTR alla NER: la pipeline AI per l'esame dei documenti",
                speaker_note="""
Passiamo alle applicazioni concrete. Vi mostrerò sei aree dove l'AI è già utilizzabile oggi nella pratica forense — non in modo sperimentale, ma con strumenti stabili, documentati e con dataset di riferimento consolidati.
""")

    # 1. HTR
    add_two_col(prs,
        "1. Trascrizione del Testo Manoscritto (HTR)",
        "Come funziona",
        [
            "HTR = Handwritten Text Recognition",
            "TrOCR: encoder visivo BEiT + decoder testuale RoBERTa",
            "L'immagine viene convertita in sequenza di token",
            "Documenti multi-riga: ogni riga viene segmentata e trascritta separatamente",
            "Output: testo digitale ricercabile e processabile",
        ],
        "Applicazioni forensi",
        [
            "Trascrizione automatica di lettere anonime o minatorie",
            "Digitalizzazione di atti giudiziari storici manoscritti",
            "Pre-elaborazione per pipeline di identificazione dell'autore",
            "Calcolo automatico del Character Error Rate (CER)",
        ],
        speaker_note="""
HTR significa Handwritten Text Recognition: il sistema legge la scrittura a mano come farebbe un dattilografo molto veloce e molto paziente. Il modello TrOCR è addestrato su milioni di esempi di testo manoscritto. Non interpreta il significato, non giudica la calligrafia — converte l'immagine in testo digitale ricercabile. Per una lettera anonima, questo significa: in pochi secondi avete un documento cercabile per parole chiave. Tecnicamente: BEiT è l'encoder visivo che 'vede' l'immagine, RoBERTa è il decoder linguistico che produce il testo.
""")

    add_bullet(prs, "1. HTR — Casi d'Uso Forensi: Esempi Concreti", [
        "Caso: lettera anonima con minacce estorsive — l'HTR produce in secondi una trascrizione completa e ricercabile",
        "Caso: testamento olografo storico in grafia difficile — l'HTR supporta la lettura e la digitalizzazione",
        "Limite: l'HTR non interpreta il contenuto, non attribuisce l'autore",
        "Il testo trascritto diventa l'input per le fasi successive della pipeline:",
        "  Ricerca di pattern linguistici",
        "  Estrazione di entità nominate (NER)",
        "  Analisi stilistica per identificazione dell'autore",
    ], note="Modello: microsoft/trocr-base-handwritten  |  Supporto: testo scritto a mano in lingue romanze",
    speaker_note="""
Un esempio concreto: una lettera minatoria scritta a mano arriva in ufficio. Il sistema trascrive il testo in secondi. Il testo trascritto viene poi passato al modulo NER per estrarre nomi, luoghi, organizzazioni — senza leggere manualmente ogni parola. Questo è il potere della pipeline: ogni strumento fa il suo pezzo, il risultato finale è molto più della somma delle parti. Il limite principale: l'HTR trascrive quello che vede, possono esserci errori su calligrafia molto personale o degradata.
""")

    # 2. Verifica firma — with siamese diagram in left column
    add_two_col(prs,
        "2. Verifica dell'Autenticità della Firma",
        "Come funziona",
        [
            "SigNet: rete siamese addestrata su dataset GPDS",
            "I due rami elaborano firma di riferimento e firma in esame",
            "Ogni firma → vettore di caratteristiche (embedding)",
            "Distanza coseno > 0.35 → classificata come FALSA",
        ],
        "Applicazioni forensi",
        [
            "Verifica di firme su assegni bancari",
            "Autenticazione di firme su contratti, testamenti, atti notarili",
            "Rilevamento di firme tracciate o riprodotte digitalmente",
            "Confronto sistematico di grandi lotti di firme",
            "Output: score numerico per il perito — non un verdetto",
        ],
        left_image=imgs.get("siamese_network"),
        speaker_note="""
La rete siamese è elegante nella sua architettura: due rami identici elaborano la firma di riferimento e quella in esame, producendo per ognuna un vettore di caratteristiche nello stesso spazio matematico. Si misura poi la distanza tra questi vettori. Se la distanza è grande, le firme sono stilisticamente diverse. Il modello è addestrato su migliaia di coppie firma-autentica/firma-falsa dal dataset GPDS. Non riconosce 'questa è la firma di Mario' — riconosce pattern di tratti, proporzioni, velocità.
""")

    add_bullet(prs, "2. Verifica Firma — Note Tecniche per l'Esaminatore", [
        "Il modello SigNet funziona meglio su firme isolate, su sfondo bianco, ad alta risoluzione",
        "Limite principale: le «skilled forgeries» (imitatori abili) possono ingannare il modello",
        "I campioni demo sono pre-selezionati dal database CEDAR: il modello rileva correttamente la contraffazione",
        "Interpretazione: più alta è la distanza coseno, più probabile è la contraffazione",
        "In casi borderline (distanza vicina alla soglia): necessaria revisione manuale dell'esaminatore",
    ], note="Modello: SigNet (luizgh/sigver) — pesi pre-addestrati su GPDS, campioni demo da CEDAR",
    speaker_note="""
Alcune precisazioni importanti per l'uso forense. Il modello funziona meglio su firme isolate, ad alta risoluzione, su sfondo chiaro. Le 'skilled forgeries' — imitatori molto abili che hanno esercitato a lungo la firma del soggetto — possono ingannare il modello, proprio come possono ingannare un esaminatore. Per questo il risultato va sempre interpretato dall'esaminatore nel contesto del caso: chi sono i possibili falsari? Avevano accesso a campioni della firma? Il sistema fornisce un dato quantitativo — il giudizio rimane umano.
""")

    # 3. Rilevamento firma — with DETR diagram
    add_two_col(prs,
        "3. Rilevamento Automatico di Firme nei Documenti",
        "Come funziona",
        [
            "Conditional DETR: Transformer encoder-decoder per rilevamento di oggetti",
            "Fine-tuned su dataset di documenti con firme annotate",
            "Produce bounding box con score di confidenza",
            "Gestisce pagine con firme multiple o parzialmente coperte",
        ],
        "Applicazioni forensi",
        [
            "Estrazione automatica di firme da contratti multi-pagina",
            "Screening di grandi archivi documentali",
            "Pipeline: Rileva → Estrai → Verifica con SigNet",
            "Report con posizione e confidenza di ogni firma",
        ],
        left_image=imgs.get("detr_detection"),
        speaker_note="""
Conditional DETR è un modello Transformer per il rilevamento di oggetti basato su ResNet-50 come backbone e un decoder Transformer per la predizione dei bounding box. Qui viene fine-tuned specificamente per trovare firme nei documenti con mAP@50 del 93,65%. La pipeline naturale è: trova le firme nel documento con Conditional DETR, poi passa ognuna a SigNet per verificarne l'autenticità. Questo è lo screening automatico di un intero contratto multi-pagina in pochi secondi.
""")

    # 4. Identificazione scrittore
    add_two_col(prs,
        "4. Identificazione dello Scrittore",
        "Come funziona",
        [
            "Feature extraction: HOG (Histogram of Oriented Gradients)",
            "Cattura stile, ritmo e distribuzione della scrittura",
            "Classificatore SVM addestrato sui campioni di riferimento",
            "Output: lista ordinata di autori candidati con score",
            "Validazione: cross-validation leave-one-out",
        ],
        "Applicazioni forensi",
        [
            "Attribuzione di lettere minatorie o estorsive anonime",
            "Verifica della paternità di documenti contestati",
            "Ricerca sulla provenienza di manoscritti storici",
            "Confronto tra campione anonimo e archivio di riferimento",
        ],
        speaker_note="""
Il sistema estrae le caratteristiche stilistiche della scrittura: HOG cattura come i tratti sono distribuiti e orientati nello spazio, LBP cattura le texture locali del tratto, le statistiche di run-length catturano la regolarità. L'SVM impara a distinguere lo stile di scrittore A da quello di scrittore B. Il risultato è una classifica di probabilità — non un verdetto di autoria, ma una lista di candidati su cui concentrare l'esame peritale. È il primo filtro quantitativo in un'indagine di attribuzione.
""")

    # 5. Analisi grafologica
    add_two_col(prs,
        "5. Analisi delle Caratteristiche Grafologiche",
        "Come funziona",
        [
            "Computer vision classica (OpenCV) per segmentazione",
            "Pipeline: binarizzazione → de-noise → raddrizzamento",
            "Misurazioni estratte:",
            "  Angolo di inclinazione delle lettere",
            "  Spaziatura tra parole e caratteri",
            "  Altezza e larghezza media dei caratteri",
            "  Pressione del tratto (intensità pixel)",
        ],
        "Applicazioni forensi",
        [
            "Supporto alla testimonianza peritale con misurazioni oggettive",
            "Confronto tra campione di riferimento e documento contestato",
            "Rilevamento di variazioni anomale: stress, camouflage, mano diversa",
            "Dashboard visiva: metriche con annotazioni sull'immagine originale",
        ],
        speaker_note="""
Questo è il modulo più vicino al lavoro tradizionale del perito grafologo: misurare le caratteristiche della scrittura. La differenza è che la macchina le misura con precisione numerica e in modo identico su ogni campione. Un perito esperto potrebbe dire 'questa scrittura ha un'inclinazione destra moderata'. Il sistema dice 'inclinazione media 12.3 gradi, deviazione standard 2.1 gradi'. Queste misurazioni sono direttamente comparabili tra campioni diversi e documentabili nel referto.
""")

    # 6. NER
    add_two_col(prs,
        "6. Riconoscimento delle Entità Nominate (NER)",
        "Come funziona",
        [
            "NER: classificazione automatica di entità nel testo",
            "WikiNEural: BERT multilingue, addestrato su Wikipedia in 9 lingue",
            "Schema BIO: B-PER, I-PER, B-ORG, B-LOC, B-MISC, O",
            "Output: lista di entità con tipo, testo e confidenza",
        ],
        "Applicazioni forensi",
        [
            "Identificare persone, luoghi, organizzazioni in lettere anonime",
            "Analizzare testamenti per estrarre nomi di eredi e beni",
            "Costruire un grafo delle relazioni tra entità",
            "Screening di archivi per la presenza di nomi specifici",
            "Pipeline completa: immagine → HTR → NER → grafo relazioni",
        ],
        speaker_note="""
NER significa Named Entity Recognition: il sistema legge il testo e identifica automaticamente le entità nominate — persone, organizzazioni, luoghi. Il modello WikiNEural è addestrato su Wikipedia in 9 lingue, quindi funziona bene su testi italiani senza configurazione aggiuntiva. È particolarmente potente come secondo passo dopo la trascrizione HTR: dall'immagine manoscritto alle entità strutturate in pochi secondi.
""")

    add_bullet(prs, "6. NER — Esempio di Applicazione Forense", [
        "Testo trascritto da HTR (testamento olografo):",
        "  «Io sottoscritto Mario Rossi, nato a Roma il 12 marzo 1950, nomino mio erede Giovanni Costa»",
        "Output NER:",
        "  [Persona] Mario Rossi          (confidenza 99%)",
        "  [Luogo]   Roma                 (confidenza 98%)",
        "  [Persona] Giovanni Costa       (confidenza 97%)",
        "Applicazione: costruzione automatica del grafo ereditario da documenti multipli",
        "Lingue supportate: italiano, inglese, tedesco, spagnolo, francese, russo, olandese, polacco, cinese",
    ], note="Modello: Babelscape/wikineural-multilingual-ner  |  Supporto multilingue nativo",
    speaker_note="""
Questo è l'esempio più concreto della pipeline completa. Un testamento olografo viene fotografato, TrOCR lo trascrive, WikiNEural estrae le entità. In pochi secondi avete una lista strutturata di tutte le persone, i luoghi e le organizzazioni menzionati nel documento — senza leggere manualmente ogni parola. Questo diventa ancora più potente quando avete decine di documenti: il sistema costruisce automaticamente un grafo delle relazioni tra tutti i soggetti citati nell'intero archivio.
""")

    # =========================================================
    # SEZIONE 4 — GraphoLab: i Laboratori
    # =========================================================
    add_section(prs, "04", "GraphoLab — I Laboratori",
                "Sette laboratori dimostrativi e un'applicazione interattiva",
                speaker_note="""
Ora vediamo concretamente GraphoLab: i sette laboratori dimostrativi e la demo interattiva. Tutto il codice è aperto, documentato e può essere eseguito localmente o via Docker senza configurazione complessa.
""")

    add_table_slide(prs,
        "GraphoLab: Panoramica dei 7 Laboratori",
        ["Lab", "Titolo", "Tecnica AI", "Modello / Strumento"],
        [
            ["01", "Introduzione", "—", "Panoramica concettuale"],
            ["02", "Riconoscimento testo manoscritto", "Transformer OCR", "TrOCR (microsoft)"],
            ["03", "Verifica autenticità firma", "Siamese Network", "SigNet (luizgh/sigver)"],
            ["04", "Rilevamento firma nei documenti", "Object Detection", "Conditional DETR (tech4humans)"],
            ["05", "Identificazione dello scrittore", "HOG + SVM", "scikit-learn"],
            ["06", "Analisi caratteristiche grafologiche", "Image Processing", "OpenCV"],
            ["07", "Riconoscimento entità nominate", "Token Classification", "WikiNEural (Babelscape)"],
        ],
        speaker_note="""
Ecco la mappa completa dei sette lab. Ognuno è autonomo — potete eseguirli nell'ordine che preferite. Il Lab 01 non ha codice: è una presentazione concettuale pensata per questo tipo di pubblico. Dal Lab 02 in poi c'è codice funzionante, dati di esempio inclusi, e una sezione dedicata ai casi d'uso forensi. I modelli scaricano automaticamente da Hugging Face al primo avvio.
""")

    add_bullet(prs, "Lab 01 — Introduzione: AI e Grafologia Forense", [
        "Notebook in stile presentazione, senza codice eseguibile",
        "Destinato a tutti i livelli — inclusi stakeholder non tecnici",
        "Contenuti:",
        "  Definizione di grafologia forense e ambiti di applicazione",
        "  Mappa concettuale della pipeline: acquisizione → preprocessing → AI → referto",
        "  Panoramica dei laboratori successivi",
    ], note="File: notebooks/01_intro_forensic_graphology.ipynb",
    speaker_note="""
Il Lab 01 è pensato proprio per chi è in sala oggi: una presentazione interattiva che spiega i concetti senza richiedere competenze di programmazione. È anche ottimo materiale da condividere con clienti, colleghi magistrati o consulenti tecnici di parte che vogliono capire cosa fa l'AI in questo contesto, senza dover eseguire codice.
""")

    add_bullet_with_image(prs, "Lab 02 — Riconoscimento del Testo Manoscritto",
        [
            "Modello: TrOCR (microsoft/trocr-base-handwritten) via Hugging Face",
            "Architettura: encoder visivo BEiT + decoder testuale RoBERTa",
            "Demo:",
            "  Trascrizione di testo manoscritto su riga singola",
            "  Documento multi-riga con segmentazione automatica",
            "  Calcolo opzionale del Character Error Rate (CER)",
        ],
        imgs.get("handwritten"),
        note="File: notebooks/02_handwritten_ocr_trocr.ipynb  |  Prerequisiti: transformers, torch, Pillow",
        font_size=17,
        speaker_note="""
TrOCR è il modello di punta di Microsoft per il riconoscimento del testo manoscritto, disponibile gratuitamente su Hugging Face. Scarica automaticamente al primo avvio — circa 350 MB. Il lab include tre demo: testo su riga singola, documento multi-riga con segmentazione automatica delle righe, e calcolo del CER per valutare l'accuratezza su un testo di riferimento noto.
""")

    add_bullet_with_image(prs, "Lab 03 — Verifica dell'Autenticità della Firma",
        [
            "Modello: SigNet con pesi pre-addestrati su dataset GPDS",
            "Demo:",
            "  Caricamento firma di riferimento e firma in esame",
            "  Estrazione embedding con il codificatore SigNet",
            "  Calcolo score di similarità coseno",
            "  Verdetto: AUTENTICA / FALSA + score di confidenza",
            "Campioni demo pre-selezionati dal database CEDAR",
        ],
        imgs.get("signatures_comparison"),
        note="File: notebooks/03_signature_verification_siamese.ipynb  |  Richiede: models/signet.pth",
        font_size=17,
        speaker_note="""
Questo è il cuore del sistema per la verifica delle firme. SigNet è il modello standard nel settore — i suoi pesi pre-addestrati sono disponibili pubblicamente nel repository luizgh/sigver. I campioni di demo provengono dal database CEDAR e sono pre-selezionati per mostrare sia casi di firma autentica che di contraffazione chiaramente rilevabile. Nella demo Gradio potete caricare le vostre firme e vedere il risultato in tempo reale.
""")

    add_bullet(prs, "Lab 04 — Rilevamento Firma nei Documenti", [
        "Modello: Conditional DETR fine-tuned (tech4humans/conditional-detr-50-signature-detector)",
        "Modello pubblico su Hugging Face (Apache 2.0) — nessun token necessario",
        "Demo:",
        "  Caricamento di un documento scansionato",
        "  Inferenza Conditional DETR: localizzazione firme con bounding box",
        "  Ritaglio automatico di ogni firma → input per Lab 03",
        "Pipeline dimostrativa: Rileva → Estrai → Verifica",
    ], note="File: notebooks/04_signature_detection_detr.ipynb  |  Prerequisiti: transformers, timm, opencv-python",
    speaker_note="""
Conditional DETR è disponibile pubblicamente su Hugging Face senza token di autenticazione, con licenza Apache 2.0. Il lab mostra la pipeline completa: il documento entra, le firme vengono rilevate e ritagliate, ogni ritaglio viene passato a SigNet per la verifica. Non è richiesta alcuna configurazione aggiuntiva oltre all'installazione di transformers e timm.
""")

    add_bullet(prs, "Lab 05 — Identificazione dello Scrittore", [
        "Tecnica: estrazione feature HOG + classificatore SVM (scikit-learn)",
        "Demo:",
        "  Caricamento di campioni di scrittori noti (dataset IAM subset)",
        "  Estrazione delle feature stilistiche da ogni campione",
        "  Addestramento SVM con cross-validation leave-one-out",
        "  Input anonimo → lista ordinata di autori candidati con score",
    ], note="File: notebooks/05_writer_identification.ipynb  |  Prerequisiti: torch, scikit-learn",
    speaker_note="""
Il lab sull'identificazione dello scrittore usa scikit-learn — nessuna GPU necessaria, gira su qualsiasi laptop. Il database di esempio include cinque scrittori con stili diversi. Per un uso forense reale, sostituite questo database con campioni autentici degli scrittori candidati: il sistema si riaddestra automaticamente. L'output è una classifica di probabilità — un punto di partenza quantitativo per l'indagine peritale.
""")

    add_bullet(prs, "Lab 06 — Analisi delle Caratteristiche Grafologiche", [
        "Strumenti: OpenCV + numpy + scipy",
        "Pipeline: binarizzazione (Otsu) → de-noise → raddrizzamento",
        "Metriche estratte e visualizzate:",
        "  Angolo di inclinazione delle lettere",
        "  Spaziatura media tra parole e caratteri",
        "  Distribuzione dell'intensità dei pixel (stima pressione tratto)",
        "Confronto opzionale affiancato tra due campioni",
    ], note="File: notebooks/06_graphological_feature_analysis.ipynb  |  Prerequisiti: opencv-python, scipy",
    speaker_note="""
L'analisi grafologica usa OpenCV — libreria di computer vision classica, stabile, documentata e gratuita. Non richiede GPU. La pipeline produce una dashboard visiva con le metriche annotate sull'immagine originale — ottimo per la presentazione in tribunale o per il referto peritale. Il confronto affiancato tra due campioni rende immediatamente visibili le differenze nelle caratteristiche misurate.
""")

    add_bullet(prs, "Lab 07 — Riconoscimento delle Entità Nominate (NER)", [
        "Modello: Babelscape/wikineural-multilingual-ner (BERT, 9 lingue)",
        "Demo 1: NER su testo italiano (testamento / dichiarazione giurata)",
        "Demo 2: NER su testo inglese (verifica supporto multilingue)",
        "Demo 3: pipeline completa HTR → NER (immagine → entità)",
        "Demo 4: distribuzione delle entità e analisi della confidenza",
        "Output: testo con evidenziazione colorata + tabella riepilogativa",
    ], note="File: notebooks/07_named_entity_recognition.ipynb  |  Prerequisiti: transformers, torch",
    speaker_note="""
Il Lab 07 è il più potente per l'analisi di documenti testuali. Il modello WikiNEural funziona su 9 lingue senza configurazione aggiuntiva — italiano, inglese, tedesco, spagnolo e altre. La Demo 3 è la più significativa: parte da un'immagine di testo manoscritto e arriva alle entità nominate estratte, tutto in un unico flusso automatico. L'output con evidenziazione colorata è direttamente utilizzabile in un referto.
""")

    add_bullet(prs, "Demo Interattiva Gradio — 6 Tab in Italiano", [
        "Applicazione web: python app/grapholab_demo.py  →  http://localhost:7860",
        "Tab 1 — OCR Manoscritto: carica immagine → testo trascritto (supporto multi-riga)",
        "Tab 2 — Verifica Firma: carica due firme → verdetto AUTENTICA / FALSA + score",
        "Tab 3 — Rilevamento Firma: carica documento → immagine annotata con bounding box",
        "Tab 4 — Riconoscimento Entità: inserisci testo → entità evidenziate + tabella",
        "Tab 5 — Identificazione Scrittore: carica campione → candidati ordinati per probabilità",
        "Tab 6 — Analisi Grafologica: carica testo manoscritto → dashboard di metriche",
        "Tutta l'interfaccia è in italiano",
    ], speaker_note="""
La demo interattiva è l'interfaccia pensata per chi non vuole usare Python o Jupyter direttamente. Si apre nel browser, è completamente in italiano, e permette di caricare immagini reali e vedere i risultati in tempo reale. Perfetta per mostrare le funzionalità a clienti, colleghi o stakeholder non tecnici. Avviamola ora per la dimostrazione live — http://localhost:7860.
""")

    add_bullet_with_image(prs,
        "Pipeline End-to-End: dal Documento al Referto Forense",
        [
            "1. Conditional DETR individua le firme nel documento → estrae le immagini",
            "2. SigNet confronta le firme con i riferimenti → verdetto autenticità",
            "3. TrOCR trascrive il testo del documento riga per riga",
            "4. WikiNEural NER estrae persone, luoghi, organizzazioni",
            "5. Analisi grafologica: inclinazione, pressione, spaziatura",
            "Output: referto con score autenticità, entità nominate, metriche grafologiche",
        ],
        imgs.get("pipeline_e2e"),
        note="La pipeline è modulare: ogni lab può essere usato autonomamente o in combinazione.",
        speaker_note="""
Questo diagramma mostra come tutti i pezzi si connettono in una pipeline end-to-end. Un documento entra nel sistema: Conditional DETR trova le firme, SigNet le verifica, TrOCR trascrive il testo, NER estrae le entità, l'analisi grafologica misura le caratteristiche. L'output finale è un referto strutturato con dati quantitativi per ogni dimensione di analisi. La pipeline è modulare: ogni componente può essere usato autonomamente o in combinazione a seconda del tipo di caso.
""")

    add_bullet(prs, "Come Avviare GraphoLab", [
        "Modalità locale (Python 3.11):",
        "  py -3.11 -m venv venv && venv\\Scripts\\activate",
        "  pip install -r requirements.txt",
        "  jupyter lab notebooks/     →  http://localhost:8888",
        "  python app/grapholab_demo.py  →  http://localhost:7860",
        "Modalità Docker (GPU supportata via WSL2):",
        "  docker compose up jupyter   →  http://localhost:8888 (token: grapholab)",
        "  docker compose up gradio    →  http://localhost:7860",
        "Cache modelli nel volume Docker grapholab-hf-cache (download unico)",
    ], font_size=17, speaker_note="""
Per chi vuole provare in locale: serve Python 3.11, poi quattro comandi e siete operativi. Per chi preferisce Docker: due comandi e tutto funziona in un ambiente isolato, senza toccare il sistema. I modelli vengono scaricati automaticamente la prima volta — Hugging Face li mette in cache nel volume Docker, quindi i download successivi sono istantanei. Nessun token HF necessario.
""")

    # =========================================================
    # SEZIONE 5 — Considerazioni Etiche e Conclusioni
    # =========================================================
    add_section(prs, "05", "Considerazioni Etiche\ne Conclusioni",
                "L'AI al servizio della giustizia: opportunità e responsabilità",
                speaker_note="""
Chiudiamo con le considerazioni etiche e legali. Questo è un argomento che i professionisti forensi conoscono bene: la tecnologia è neutrale, ma l'uso che se ne fa non lo è. Alcune riflessioni fondamentali prima di concludere.
""")

    add_bullet(prs, "Bias nei Dati e Spiegabilità", [
        "Bias nei dati di addestramento:",
        "  I modelli riflettono i dati su cui sono stati addestrati",
        "  Soluzione: validare il modello sul tipo specifico di documenti del caso",
        "Spiegabilità in tribunale:",
        "  Le conclusioni peritali devono essere giustificabili e aperte alla contestazione",
        "  La soglia decisionale deve essere documentata e motivata",
        "  L'AI va usata in modo spiegabile in termini non tecnici",
    ], speaker_note="""
Due punti critici. Il bias: un modello addestrato solo su firme europee di adulti sani in condizioni normali potrebbe non funzionare bene su firme di anziani, di persone con condizioni neurologiche, o di culture grafiche diverse. La validazione sul tipo specifico di documenti del caso è quindi fondamentale prima dell'uso forense. La spiegabilità: un giudice deve poter capire come il sistema ha prodotto la sua conclusione. 'Il modello di deep learning ha detto così' non è una risposta accettabile. L'esaminatore deve saper spiegare la metodologia in termini comprensibili.
""")

    add_bullet(prs, "Catena di Custodia e Standard di Riferimento", [
        "Catena di custodia digitale:",
        "  I file originali devono essere conservati inalterati (hash crittografico)",
        "  Ogni passo dell'analisi AI deve essere documentato: versione modello, parametri, data/ora",
        "Standard internazionali di riferimento:",
        "  OSAC — Organization of Scientific Area Committees for Forensic Science",
        "  SWGDOC — Scientific Working Group for Questioned Documents",
        "  Entrambi stanno sviluppando linee guida per gli strumenti computazionali nell'esame dei documenti",
    ], speaker_note="""
La catena di custodia digitale è fondamentale e spesso sottovalutata. I file originali devono essere conservati con hash crittografico — SHA-256 o simili — per dimostrare che non sono stati alterati tra l'acquisizione e l'analisi. Ogni passo dell'analisi AI deve essere documentato: quale versione del modello, quali parametri, quando è stata eseguita. OSAC e SWGDOC sono gli organismi di riferimento internazionale per gli standard forensi e stanno sviluppando linee guida specifiche per gli strumenti computazionali nell'esame dei documenti questionati.
""")

    add_bullet(prs, "AI come Strumento, Non come Sostituto", [
        "L'AI non sostituisce il perito calligrafo — lo potenzia",
        "Cosa rimane esclusivamente umano:",
        "  L'interpretazione dei risultati nel contesto del caso specifico",
        "  La conoscenza specializzata delle dinamiche della scrittura",
        "  Il giudizio professionale e la responsabilità legale del referto",
        "L'AI comprime in minuti ciò che richiederebbe giorni, trasformando impressioni visive in misurazioni precise",
        "Il risultato è un'analisi più rigorosa, più difendibile e più utile alla ricerca della verità",
    ], speaker_note="""
Il messaggio conclusivo e più importante: l'AI è uno strumento nelle mani del professionista. Il referto forense è firmato dall'esperto, non dal computer. La responsabilità legale è dell'esaminatore, non dell'algoritmo. L'AI comprime in minuti ciò che richiederebbe giorni — ma l'interpretazione, il giudizio contestuale, la responsabilità rimangono esclusivamente umani. L'obiettivo di GraphoLab è fornire questi strumenti in modo accessibile, documentato e utilizzabile nella pratica professionale quotidiana.
""")

    # Final slide
    slide = _blank_slide(prs)
    _fill_bg(slide, NAVY)
    _add_rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55), GOLD)
    _add_textbox(slide, "Grazie per l'attenzione",
                 Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
                 font_name=TITLE_FONT, font_size=40, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Domande e discussione",
                 Inches(1), Inches(3.2), Inches(11.3), Inches(0.8),
                 font_size=26, color=GOLD, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "GraphoLab — Intelligenza Artificiale in Grafologia Forense",
                 Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
                 font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.CENTER)
    _set_speaker_note(slide, """
Grazie per l'attenzione. Sono a disposizione per qualsiasi domanda — sia sui contenuti tecnici che sulle applicazioni pratiche nel vostro lavoro quotidiano. Se volete sperimentare con i vostri dati, i lab sono disponibili e possiamo pianificare una sessione pratica. Il repository GraphoLab è pubblico su GitHub — trovate il link nella slide.
""")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("Generazione diagrammi matplotlib...")
    imgs = generate_diagrams()

    print(f"\nCreazione presentazione ({len(imgs)} figure disponibili)...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build(prs, imgs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"\nPresentazione salvata: {OUTPUT}")
    print(f"Numero di slide: {len(prs.slides)}")


if __name__ == "__main__":
    main()
